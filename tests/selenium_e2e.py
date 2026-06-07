"""End-to-end test for PDF Merger Pro.

Runs a real (headless-new) Chrome — where pdf.js canvas rendering works — serves
the app from the project root, uploads real PDF/PNG/DOCX/PPTX files and exercises
the full UI: thumbnails, PDF text recognition + editing, image cropping, preview
and a mixed merge (verifying the downloaded PDF).

Setup:
    pip install selenium reportlab python-docx python-pptx Pillow pypdf
    (Chrome must be installed; Selenium Manager fetches the matching driver.)

Run:
    python tests/selenium_e2e.py
Exit code 0 means all checks passed.
"""
import os, sys, time, tempfile, threading, http.server, socketserver, glob
try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

PROJECT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
tmp = tempfile.mkdtemp(prefix="pdfmerger_sel_")
dl_dir = os.path.join(tmp, "downloads"); os.makedirs(dl_dir, exist_ok=True)
results = []
def log(name, ok, detail=""):
    results.append((name, ok, detail))
    print(("PASS" if ok else "FAIL"), "-", name, (":: " + detail) if detail else "")

# ---------- 1. create real test files ----------
from reportlab.pdfgen import canvas as rcanvas
pdf_path = os.path.join(tmp, "test-doc.pdf")
c = rcanvas.Canvas(pdf_path, pagesize=(400, 500))
c.setFont("Helvetica", 24); c.drawString(60, 430, "Vertraulich")
c.setFont("Helvetica", 13); c.drawString(60, 400, "Zeile zwei zum Testen")
c.showPage()
c.setFont("Helvetica", 18); c.drawString(60, 450, "Seite zwei")
c.showPage(); c.save()

from PIL import Image, ImageDraw
png_path = os.path.join(tmp, "photo.png")
im = Image.new("RGB", (320, 200), (59, 130, 246))
ImageDraw.Draw(im).rectangle([40, 40, 280, 160], fill=(255, 255, 255))
im.save(png_path)

import docx
docx_path = os.path.join(tmp, "report.docx")
d = docx.Document(); d.add_heading("Test Report", 0)
d.add_paragraph("Erster Absatz im Word-Dokument."); d.add_paragraph("Zweiter Absatz.")
d.save(docx_path)

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.shapes import MSO_SHAPE
pptx_path = os.path.join(tmp, "slides.pptx")
prs = Presentation()
s = prs.slides.add_slide(prs.slide_layouts[0])
s.shapes.title.text = "Selenium Deck"; s.placeholders[1].text = "Untertitel"
# slide 2 with a colored shape — regression guard: html2canvas dropped SVG shapes
s2 = prs.slides.add_slide(prs.slide_layouts[6])
box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1), Inches(4), Inches(2))
box.fill.solid(); box.fill.fore_color.rgb = PRGB(0x3B, 0x82, 0xF6); box.text_frame.text = "Phase 1"
prs.save(pptx_path)
print("test files in", tmp)

# ---------- 2. start static server from project ----------
os.chdir(PROJECT)
class Q(http.server.SimpleHTTPRequestHandler):
    def log_message(self, *a): pass
httpd = socketserver.TCPServer(("127.0.0.1", 0), Q)
port = httpd.server_address[1]
threading.Thread(target=httpd.serve_forever, daemon=True).start()
base = f"http://127.0.0.1:{port}/index.html"
print("server", base)

# ---------- 3. selenium ----------
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait

opts = Options()
opts.add_argument("--headless=new")
opts.add_argument("--window-size=1500,1100")
opts.add_argument("--no-sandbox")
opts.add_argument("--disable-gpu")
opts.add_experimental_option("prefs", {
    "download.default_directory": dl_dir,
    "download.prompt_for_download": False,
    "savefile.default_directory": dl_dir,
})
driver = webdriver.Chrome(options=opts)
driver.set_page_load_timeout(60)
W = WebDriverWait(driver, 30)

def js(s): return driver.execute_script(s)

try:
    driver.get(base)
    W.until(lambda dr: dr.execute_script(
        "return document.readyState==='complete' && typeof handleFiles==='function' "
        "&& typeof PDFLib==='object' && typeof Cropper==='function' && typeof pdfjsLib==='object' "
        "&& typeof window.docx==='object' && typeof html2canvas==='function' "
        "&& !!(window.jQuery&&window.jQuery.fn.pptxToHtml)"))
    log("app + all libraries loaded", True, js("return document.title"))
    log("visibilityState is 'visible' (needed for pdf.js)", js("return document.visibilityState") == "visible", js("return document.visibilityState"))

    inp = driver.find_element(By.ID, "fileInput")
    inp.send_keys("\n".join([pdf_path, png_path, docx_path, pptx_path]))

    W.until(lambda dr: dr.execute_script("return document.querySelectorAll('#resultsList .file-item').length") == 4)
    time.sleep(2.5)
    log("4 files added", js("return document.querySelectorAll('#resultsList .file-item').length") == 4)

    metas = js("return [...document.querySelectorAll('#resultsList .file-meta')].map(e=>e.textContent)")
    log("row metadata", True, " || ".join(metas))

    thumbs = js("return [...document.querySelectorAll('#resultsList .file-item')].map(it=>{"
                "const img=it.querySelector('.thumb-img'); return img && !img.classList.contains('hidden') "
                "&& (img.src||'').startsWith('data:')?1:0;})")
    log("PDF thumbnail rendered (pdf.js)", thumbs[0] == 1, f"thumb[0]={thumbs[0]}")
    log("image thumbnail", thumbs[1] == 1, f"thumb[1]={thumbs[1]}")
    log("DOCX thumbnail rendered", thumbs[2] == 1, f"thumb[2]={thumbs[2]}")
    log("PPTX thumbnail rendered (pptxjs + modern-screenshot)", thumbs[3] == 1, f"thumb[3]={thumbs[3]}")
    log("totals shown", True, js("return document.getElementById('totalPageInfo').textContent"))

    # PDF text editor
    driver.find_element(By.CSS_SELECTOR, "#resultsList .file-item:nth-child(1) .edit-btn").click()
    W.until(lambda dr: dr.execute_script("return document.getElementById('editModal').classList.contains('active')"))
    W.until(lambda dr: dr.execute_script("return document.querySelectorAll('#editOverlay .anno-recognized').length") > 0)
    recog = js("return document.querySelectorAll('#editOverlay .anno-recognized').length")
    sample = js("return (document.querySelector('#editOverlay .anno-recognized')||{}).textContent || ''")
    log("PDF text recognized as editable boxes", recog > 0, f"boxes={recog} sample='{sample}'")
    js("var el=document.querySelector('#editOverlay .anno-recognized'); el.focus(); "
       "el.textContent='GEAENDERT'; el.dispatchEvent(new Event('input',{bubbles:true}));")
    log("changed box flagged", js("return document.querySelector('#editOverlay .anno-recognized').classList.contains('anno-changed')"))
    driver.find_element(By.ID, "editApply").click()
    W.until(lambda dr: not dr.execute_script("return document.getElementById('editModal').classList.contains('active')"))
    m = js("return document.querySelector('#resultsList .file-item:nth-child(1) .file-meta').textContent")
    log("PDF marked edited", "bearbeitet" in m or "edited" in m, m)

    # crop image
    driver.find_element(By.CSS_SELECTOR, "#resultsList .file-item:nth-child(2) .crop-btn").click()
    W.until(lambda dr: dr.execute_script("return document.getElementById('cropModal').classList.contains('active')"))
    W.until(lambda dr: dr.execute_script("return !!document.querySelector('.cropper-container')"))
    log("crop modal + cropper.js init", True)
    js("if(cropper) cropper.setData({x:60,y:40,width:120,height:80});")
    driver.find_element(By.ID, "cropApply").click()
    W.until(lambda dr: not dr.execute_script("return document.getElementById('cropModal').classList.contains('active')"))
    m = js("return document.querySelector('#resultsList .file-item:nth-child(2) .file-meta').textContent")
    log("image cropped", "zugeschnitten" in m or "cropped" in m, m)

    # preview docx
    driver.find_element(By.CSS_SELECTOR, "#resultsList .file-item:nth-child(3) .preview-btn").click()
    W.until(lambda dr: dr.execute_script("return document.getElementById('previewModal').classList.contains('active')"))
    time.sleep(0.5)
    log("DOCX preview shows page image",
        js("return !document.getElementById('previewImage').classList.contains('hidden') && (document.getElementById('previewImage').src||'').length>100"))
    driver.find_element(By.ID, "previewClose").click()
    time.sleep(0.3)

    # PPTX shape slide must actually render its colored shapes (regression guard)
    driver.find_element(By.CSS_SELECTOR, "#resultsList .file-item:nth-child(4) .preview-btn").click()
    W.until(lambda dr: dr.execute_script("return document.getElementById('previewModal').classList.contains('active')"))
    time.sleep(0.8)
    js("document.getElementById('pdfNext').click()")  # go to slide 2 (the shape slide)
    time.sleep(0.8)
    colored = driver.execute_async_script("""
      const cb=arguments[arguments.length-1];
      const img=new Image();
      img.onload=()=>{const c=document.createElement('canvas');c.width=img.naturalWidth;c.height=img.naturalHeight;
        const x=c.getContext('2d');x.drawImage(img,0,0);const d=x.getImageData(0,0,c.width,c.height).data;let n=0;
        for(let i=0;i<d.length;i+=4){if(Math.abs(d[i]-d[i+1])>25||Math.abs(d[i+1]-d[i+2])>25)n++;}cb(n);};
      img.onerror=()=>cb(-1);
      img.src=document.getElementById('previewImage').src;
    """)
    log("PPTX shape slide renders colored shapes (modern-screenshot/SVG)", colored > 1000, f"colored_px={colored}")
    driver.find_element(By.ID, "previewClose").click()
    time.sleep(0.3)

    # merge
    js("document.getElementById('outputFilename').value='selenium-merged';")
    driver.find_element(By.ID, "mergeBtn").click()
    deadline = time.time() + 40
    merged = None
    while time.time() < deadline:
        files = [f for f in glob.glob(os.path.join(dl_dir, "*.pdf")) if not f.endswith(".crdownload")]
        if files:
            merged = files[0]; break
        time.sleep(0.5)
    if merged:
        from pypdf import PdfReader
        reader = PdfReader(merged)
        npages = len(reader.pages)
        alltext = ""
        for p in reader.pages:
            try: alltext += (p.extract_text() or "")
            except Exception: pass
        log("merged PDF downloaded", True, f"{os.path.basename(merged)} {os.path.getsize(merged)}B {npages} pages")
        # PDF (2) + image (1) + DOCX (1) + PPTX (2) = 6
        log("merged page count = 6", npages == 6, f"pages={npages}")
        log("edited PDF text 'GEAENDERT' present in output", "GEAENDERT" in alltext)
        log("DOCX content (rasterized -> no selectable text expected)", True, f"chars_extracted={len(alltext)}")
    else:
        log("merged PDF downloaded", False, "no file in download dir")

    logs = driver.get_log("browser")
    severe = [l for l in logs if l["level"] == "SEVERE" and "tailwind" not in l["message"].lower()]
    log("no severe console errors", len(severe) == 0, (severe[0]["message"][:160] if severe else ""))

finally:
    driver.quit()
    httpd.shutdown()

passed = sum(1 for _, ok, _ in results if ok)
print("\n==== SUMMARY:", passed, "/", len(results), "passed ====")
sys.exit(0 if passed == len(results) else 1)
